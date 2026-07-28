"""Build RoadService presentation PPTX in repo root.

Run from repo root: python scripts/make_presentation.py
"""

from __future__ import annotations

import shutil
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
SHOTS = ROOT / "docs" / "screenshots"
OUT = ROOT / "RoadService_Presentation.pptx"

# Prefer Cursor-generated assets, then docs/screenshots
ASSET_CANDIDATES = [
    Path(r"C:\Users\surti\.cursor\projects\d-RoadService\assets"),
    ROOT / "docs" / "screenshots",
]

SLIDES = [
    ("shot-landing.png", "Landing page", "Public entry point with brand-first hero and sign-in CTAs."),
    ("shot-dashboard.png", "Dashboard", "Live KPIs for open, in-progress, verification, and closed issues."),
    ("shot-issues.png", "Issues + photo stages", "Surveyor before, contractor submit, and final closed photos."),
    ("shot-notifications.png", "Notifications", "Assign, submit, reject, and 24h verification warnings."),
    ("shot-mobile.png", "Mobile contractor flow", "One-tap Start / Submit / View comments + offline capture sync."),
]


def find_image(name: str) -> Path | None:
    for base in ASSET_CANDIDATES:
        p = base / name
        if p.exists():
            return p
    return None


def style_title(shape, size=32):
    for p in shape.text_frame.paragraphs:
        p.font.bold = True
        p.font.size = Pt(size)
        p.font.color.rgb = RGBColor(0xE8, 0xEE, 0xF6)
        p.font.name = "Calibri"


def add_dark_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0x0A, 0x0C, 0x10)


def main() -> None:
    SHOTS.mkdir(parents=True, exist_ok=True)
    for name, _, _ in SLIDES:
        src = find_image(name)
        if src:
            shutil.copy2(src, SHOTS / name)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Title
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_bg(slide)
    box = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(11.5), Inches(1.2))
    box.text_frame.text = "RoadService"
    style_title(box, 54)
    sub = slide.shapes.add_textbox(Inches(0.8), Inches(3.5), Inches(11.5), Inches(1.5))
    tf = sub.text_frame
    tf.text = "Road Issue Management System"
    p = tf.paragraphs[0]
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(0x3B, 0x9E, 0xFF)
    p2 = tf.add_paragraph()
    p2.text = "Survey → Repair → Verify → Close  |  Web + Mobile  |  Photo + GPS proof"
    p2.font.size = Pt(16)
    p2.font.color.rgb = RGBColor(0x8B, 0x9B, 0xB0)

    # Agenda
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_bg(slide)
    title = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11), Inches(0.8))
    title.text_frame.text = "What we built"
    style_title(title, 36)
    body = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.5), Inches(5))
    tf = body.text_frame
    bullets = [
        "Multi-role workflow: Surveyor, Contractor, Admin, Government (view-only)",
        "Photo evidence at every stage: before → after submit → final closed",
        "Notifications: assign, submit ready, reject/rework, 24h verification warning",
        "One-tap contractor actions from list + notification deep-links",
        "Offline mobile capture with later sync when network returns",
        "Map, analytics, Excel/PDF reports for oversight",
    ]
    for i, line in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"•  {line}"
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(0xE8, 0xEE, 0xF6)
        p.level = 0
        p.space_after = Pt(12)

    # Screenshot slides
    for name, heading, caption in SLIDES:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_dark_bg(slide)
        t = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12), Inches(0.55))
        t.text_frame.text = heading
        style_title(t, 28)
        c = slide.shapes.add_textbox(Inches(0.5), Inches(0.75), Inches(12), Inches(0.4))
        cp = c.text_frame.paragraphs[0]
        cp.text = caption
        cp.font.size = Pt(14)
        cp.font.color.rgb = RGBColor(0x8B, 0x9B, 0xB0)
        img = SHOTS / name
        if img.exists():
            slide.shapes.add_picture(str(img), Inches(0.7), Inches(1.25), width=Inches(12))
        else:
            miss = slide.shapes.add_textbox(Inches(0.7), Inches(3), Inches(11), Inches(1))
            miss.text_frame.text = f"Screenshot missing: {name}"
            style_title(miss, 20)

    # Workflow
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_bg(slide)
    title = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11), Inches(0.8))
    title.text_frame.text = "End-to-end process"
    style_title(title, 36)
    steps = [
        ("1. Report", "Surveyor captures defect + GPS + before photo (mobile)"),
        ("2. Assign", "Contractor gets notification and starts work in one tap"),
        ("3. Repair", "Contractor submits camera proof + remarks (offline-safe)"),
        ("4. Verify", "Surveyor/Admin approve or reject within 24 hours"),
        ("5. Close", "Final verification photo stored; Government sees full trail"),
    ]
    body = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.5), Inches(5.2))
    tf = body.text_frame
    for i, (h, d) in enumerate(steps):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = h
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0x3B, 0x9E, 0xFF)
        p2 = tf.add_paragraph()
        p2.text = d
        p2.font.size = Pt(16)
        p2.font.color.rgb = RGBColor(0xE8, 0xEE, 0xF6)
        p2.space_after = Pt(14)

    # Closing
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_bg(slide)
    box = slide.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(11.5), Inches(2))
    tf = box.text_frame
    tf.text = "Accountability with proof"
    style_title(box, 40)
    p = tf.add_paragraph()
    p.text = "Every road issue has owners, timestamps, GPS, and photos — from survey to close."
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(0x8B, 0x9B, 0xB0)
    p.alignment = PP_ALIGN.LEFT

    prs.save(OUT)
    print(f"Wrote {OUT}")


if __name__ == "__main__":
    main()
